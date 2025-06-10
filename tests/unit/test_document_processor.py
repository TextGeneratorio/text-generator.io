import requests
import pytest

from questions.document_processor import convert_to_markdown, convert_documents_async

pytestmark = [pytest.mark.integration, pytest.mark.internet]

DOCX_URL = "https://calibre-ebook.com/downloads/demos/demo.docx"
PPTX_URL = "https://filesamples.com/samples/document/pptx/sample3.pptx"
PDF_URL = "https://github.com/mozilla/pdf.js/raw/master/web/compressed.tracemonkey-pldi-09.pdf"


def download(url, path):
    resp = requests.get(url, timeout=30)
    resp.raise_for_status()
    with open(path, "wb") as f:
        f.write(resp.content)


@pytest.fixture(scope="session")
def docx_file(tmp_path_factory):
    from docx import Document
    path = tmp_path_factory.mktemp("docs") / "sample.docx"
    doc = Document()
    doc.add_paragraph("Hello World")
    doc.save(path)
    return path

@pytest.fixture(scope="session")
def pptx_file(tmp_path_factory):
    from pptx import Presentation
    path = tmp_path_factory.mktemp("docs") / "sample.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(path)
    return path

@pytest.fixture(scope="session")
def pdf_file(tmp_path_factory):
    path = tmp_path_factory.mktemp("docs") / "sample.pdf"
    download(PDF_URL, path)
    return path


def test_docx_local(docx_file):
    md = convert_to_markdown(str(docx_file))
    assert md.strip()


def test_pptx_local(pptx_file):
    md = convert_to_markdown(str(pptx_file))
    assert md.strip()


def test_pdf_local(pdf_file):
    md = convert_to_markdown(str(pdf_file))
    assert md.strip()


@pytest.mark.asyncio
async def test_convert_documents_async(docx_file, pptx_file):
    results = await convert_documents_async([str(docx_file), str(pptx_file)])
    assert len(results) == 2
    for text in results.values():
        assert text.strip()
